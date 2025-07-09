"""
Document Processor Module

Handles text extraction from various document formats (PDF, PPTX)
and text chunking for optimal retrieval performance.
"""

import io
from typing import List

from src.logging.logger import get_logger

logger = get_logger(__name__)

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file content."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=file_content, filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text.strip()
    except ImportError:
        # Fallback to pdfplumber if PyMuPDF is not available
        try:
            import pdfplumber
            
            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                text = ""
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                return text.strip()
        except ImportError:
            raise Exception("No PDF parsing library available. Please install required dependencies:\n" +
                          "pip install PyMuPDF pdfplumber\n" +
                          "or\n" +
                          "python -m pip install PyMuPDF pdfplumber")
    except Exception as e:
        raise Exception(f"Failed to extract text from PDF: {str(e)}")

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file content."""
    try:
        from pptx import Presentation
        
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text.strip()
    except ImportError:
        raise Exception("python-pptx library not available. Please install required dependency:\n" +
                      "pip install python-pptx\n" +
                      "or\n" +
                      "python -m pip install python-pptx")
    except Exception as e:
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")

def chunk_text(text: str, chunk_size: int = 350, overlap: int = 100) -> List[str]:
    """
    Split text into overlapping chunks for better retrieval.
    
    Args:
        text: The input text to chunk
        chunk_size: Number of words per chunk
        overlap: Number of words to overlap between chunks
        
    Returns:
        List of text chunks
    """
    words = text.split()
    chunks = []
    
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk.strip())
    
    return chunks

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file content."""
    try:
        from docx import Document
        
        doc = Document(io.BytesIO(file_content))
        text = ""
        
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        return text.strip()
    except ImportError:
        raise Exception("python-docx library not available. Please install required dependency:\n" +
                      "pip install python-docx\n" +
                      "or\n" +
                      "python -m pip install python-docx")
    except Exception as e:
        raise Exception(f"Failed to extract text from DOCX: {str(e)}")

def extract_text_from_txt(file_content: bytes) -> str:
    """Extract text from plain text file."""
    try:
        # Try UTF-8 first
        text = file_content.decode('utf-8')
        return text.strip()
    except UnicodeDecodeError:
        try:
            # Fallback to latin-1
            text = file_content.decode('latin-1')
            return text.strip()
        except UnicodeDecodeError:
            # Last resort - ignore errors
            text = file_content.decode('utf-8', errors='ignore')
            return text.strip()

def get_document_extractor(file_extension: str):
    """
    Get the appropriate text extractor function for a given file extension.
    
    Args:
        file_extension: File extension (e.g., '.pdf', '.pptx')
        
    Returns:
        Callable text extractor function
    """
    extractors = {
        '.pdf': extract_text_from_pdf,
        '.pptx': extract_text_from_pptx,
        '.docx': extract_text_from_docx,
        '.txt': extract_text_from_txt
    }
    
    return extractors.get(file_extension.lower())

def is_supported_document(file_extension: str) -> bool:
    """Check if a file extension is supported for document processing."""
    supported_extensions = ['.pdf', '.pptx', '.docx', '.txt']
    return file_extension.lower() in supported_extensions 